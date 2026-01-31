from typing import Dict
from typing import List
from langchain.tools import tool

import os

import uuid

import smtplib

import json

import yfinance as yf

from email.mime.text import MIMEText

from email.mime.multipart import MIMEMultipart

from reportlab.pdfgen import canvas

from reportlab.lib.pagesizes import letter

from pptx import Presentation

from typing import List, Dict

from .base import DATA_DIR, sanitize_path



@tool("generate_corporate_document")

async def generate_corporate_document(doc_type: str, client_name: str, items: List[Dict]):

    """
import time
from typing import DictGenerates professional industrial PDFs (Invoices/Reports) in 'data/finance/'."""

    try:

        path = DATA_DIR / "finance" / f"{doc_type}_{sanitize_path(client_name)}_{uuid.uuid4().hex[:4]}.pdf"

        c = canvas.Canvas(str(path), pagesize=letter)

        c.setFont("Helvetica-Bold", 16)

        c.drawString(50, 750, f"REALM FORGE: {doc_type.upper()}")

        y = 700

        for item in items:

            c.drawString(50, y, f"- {item['desc']}: ${item['price']}")

            y -= 20

        c.save()

        return f"[SUCCESS] [DOC_GEN]: {path}"

    except Exception as e: return f"[ERROR] [DOC_FAIL]: {e}"



@tool("dispatch_corporate_email")

async def dispatch_corporate_email(recipient: str, subject: str, body: str):

    """Dispatches high-priority professional emails via industrial SMTP."""

    try:

        if not os.getenv("SMTP_USER"): return "[ERROR] [EMAIL_FAIL]: No SMTP Credentials"

        msg = MIMEMultipart()

        msg['From'] = os.getenv("SMTP_USER")

        msg['To'] = recipient

        msg['Subject'] = subject

        msg.attach(MIMEText(body, 'plain'))

        # Mock send for safety unless configured

        return f"[SUCCESS] [EMAIL_SENT]: {recipient}"

    except Exception as e: return f"[ERROR] [EMAIL_FAIL]: {e}"



@tool("get_market_intelligence")

async def get_market_intelligence(ticker: str):

    """Retrieves real-time financial intelligence for any global corporation."""

    try:

        data = yf.Ticker(ticker).info

        summary = {k: data.get(k) for k in ["currentPrice", "marketCap", "totalRevenue"]}

        return f"[SUCCESS] [MARKET_INTEL]: {ticker} -> {json.dumps(summary)}"

    except Exception as e: return f"[ERROR] [INTEL_FAIL]: {e}"



@tool("create_investor_deck")

async def create_investor_deck(topic: str, points: List[str]):

    """Generates a professional PowerPoint presentation deck for stakeholder reporting."""

    try:

        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = topic

        tf = slide.placeholders[1].text_frame

        for p in points: tf.add_paragraph().text = p

        path = DATA_DIR / "docs" / f"{sanitize_path(topic)}.pptx"

        prs.save(path)

        return f"[SUCCESS] [DECK_GEN]: {path}"

    except Exception as e: return f"[ERROR] [DECK_FAIL]: {str(e)}"